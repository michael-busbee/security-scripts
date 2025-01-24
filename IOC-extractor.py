from pathlib import Path
import docx
import PyPDF2
from email import policy
from email.parser import BytesParser
from openpyxl import load_workbook  # For .xlsx files
from pptx import Presentation  # For .pptx files
from pyrtf_ng import Rtf15Reader, plaintext  # For .rtf files
import re

def read_text(file_path):

    file_path = Path(file_path)
    
    if not file_path.is_file():
        raise ValueError(f"The path '{file_path}' is not a valid file.")

    file_extension = file_path.suffix.lower()
    extracted_text = ""

    try:


        if file_extension == ".pdf":
            # Read PDF file
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                extracted_text = "".join([page.extract_text() for page in pdf_reader.pages])


        elif file_extension == ".docx":
            # Read DOCX file
            doc = docx.Document(file_path)
            extracted_text = "\n".join([paragraph.text for paragraph in doc.paragraphs])


        elif file_extension == ".eml":
            # Read EML file
            with open(file_path, 'rb') as file:
                msg = BytesParser(policy=policy.default).parse(file)
                extracted_text = msg.get_body(preferencelist=('plain')).get_content() if msg.is_multipart() else msg.get_content()
        
        
        elif file_extension == ".xlsx":
            # Read XLSX file
            workbook = load_workbook(file_path, read_only=True)
            lines = []
            for sheet in workbook:
                for row in sheet.iter_rows(values_only=True):
                    cells = [(str(cell) if cell is not None else "") for cell in row]
                    lines.append("\t".join(cells))
            extracted_text = "\n".join(lines)

        
        elif file_extension == ".pptx":
            # Read PPTX file
            presentation = Presentation(file_path)
            slides_text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slides_text.append(shape.text)
            extracted_text = "\n".join(slides_text)

        
        elif file_extension == ".rtf":
            # Read RTF file
            with open(file_path, "rb") as f:
                doc = Rtf15Reader.read(f)
            extracted_text = plaintext.PlaintextWriter.write(doc).getvalue()
        
        
        else:
            # Read text file
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                extracted_text = file.read()

    except Exception as e:
        raise ValueError(f"An error occurred while reading the file '{file_path}': {e}")

    return extracted_text



def scan_text(text):
    # Dictionary to store found IOCs
    iocs_found = {
        "ip_addresses": [],
        "domains": [],
        "urls": [],
        "dns_records": [],
        "user_agents": [],
        "md5_hashes": [],
        "sha256_hashes": [],
        "file_names": [],
        "windows_registry_keys": [],
        "email_addresses": [],
        "bitcoin_addresses": [],
        "ethereum_addresses": [],
        "twitter_handles": [],
        "telegram_usernames": [],
        "discord_usernames": [],
        "mac_addresses": []
    }

    # IP (IPv4) addresses
    iocs_found["ip_addresses"] = re.findall(
        r"\b(?:[0-9]{1,3}\.){3}[0-9]{1,3}\b", text
    )

    # Domain names (very simplified pattern)
    iocs_found["domains"] = re.findall(
        r"\b(?:[a-zA-Z0-9-]+\.)+[a-zA-Z]{2,}\b", text
    )

    # URLs (simple pattern for http/https)
    iocs_found["urls"] = re.findall(
        r"(https?://[^\s]+)", text
    )

    # DNS records (basic match for common record types)
    iocs_found["dns_records"] = re.findall(
        r"\b(?:A|AAAA|CNAME|MX|NS|SOA|TXT|SRV)\b", text, flags=re.IGNORECASE
    )

    # User-Agents (capture text after "User-Agent:" on a line)
    iocs_found["user_agents"] = re.findall(
        r"(?im)^(?:user-agent):\s*(.+)$", text
    )

    # MD5 hashes (32 hex chars)
    iocs_found["md5_hashes"] = re.findall(
        r"\b[a-fA-F0-9]{32}\b", text
    )

    # SHA256 hashes (64 hex chars)
    iocs_found["sha256_hashes"] = re.findall(
        r"\b[a-fA-F0-9]{64}\b", text
    )

    # File names (basic pattern: some chars + extension)
    iocs_found["file_names"] = re.findall(
        r"\b[\w,\s-]+\.[A-Za-z0-9]{1,5}\b", text
    )

    # Windows Registry Keys (common root keys + path)
    iocs_found["windows_registry_keys"] = re.findall(
        r"\b(HKEY_LOCAL_MACHINE|HKEY_CURRENT_USER|HKEY_CLASSES_ROOT|HKEY_USERS|HKEY_CURRENT_CONFIG)\\[^\s]*",
        text
    )

    # Email addresses
    iocs_found["email_addresses"] = re.findall(
        r"[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+", text
    )

    # Bitcoin addresses (simple pattern matching 1/3/bc1)
    iocs_found["bitcoin_addresses"] = re.findall(
        r"\b(?:[13][a-km-zA-HJ-NP-Z1-9]{25,34}|bc1[a-zA-HJ-NP-Z0-9]{39,59})\b",
        text
    )

    # Ethereum addresses (starts with 0x + 40 hex)
    iocs_found["ethereum_addresses"] = re.findall(
        r"\b0x[a-fA-F0-9]{40}\b", text
    )

    # Twitter handles (up to 15 alphanumeric or underscore chars after @)
    iocs_found["twitter_handles"] = re.findall(
        r"(?<!\S)@([A-Za-z0-9_]{1,15})(?!\S)", text
    )

    # Telegram username (5 to 32 alphanumeric/underscore chars after @)
    iocs_found["telegram_usernames"] = re.findall(
        r"(?<!\S)@[a-zA-Z0-9_]{5,32}(?!\S)", text
    )

    # Discord usernames (2 to 32 word chars + # + 4 digits)
    iocs_found["discord_usernames"] = re.findall(
        r"\b[A-Za-z0-9_]{2,32}#[0-9]{4}\b", text
    )

    # MAC addresses (6 pairs of hex separated by : or -)
    iocs_found["mac_addresses"] = re.findall(
        r"\b(?:[0-9A-Fa-f]{2}[-:]){5}[0-9A-Fa-f]{2}\b", text
    )

    return iocs_found


def main():
    print(read_text("testdoc.docx"))


if __name__ == "__main__":
    main()